#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
HVAC 관리 시스템 사용자 가이드 프레젠테이션 생성
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def create_hvac_presentation():
    """HVAC 관리 시스템 프레젠테이션 생성"""
    
    # 프레젠테이션 객체 생성
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # 색상 정의
    PRIMARY_COLOR = RGBColor(103, 126, 234)  # #677eea
    SECONDARY_COLOR = RGBColor(99, 102, 241)  # #6366f1
    SUCCESS_COLOR = RGBColor(16, 185, 129)  # #10b981
    WARNING_COLOR = RGBColor(245, 158, 11)  # #f59e0b
    DANGER_COLOR = RGBColor(239, 68, 68)  # #ef4444
    DARK_COLOR = RGBColor(51, 51, 51)  # #333333
    LIGHT_COLOR = RGBColor(248, 249, 250)  # #f8f9fa
    
    # 슬라이드 1: 표지
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 배경색
    background = slide1.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = PRIMARY_COLOR
    
    # 제목
    title_box = slide1.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "HVAC 관리 시스템"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER
    
    # 부제목
    subtitle_box = slide1.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "사용자별 운영 가이드"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(32)
    subtitle_para.font.color.rgb = RGBColor(255, 255, 255)
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # 날짜
    date_box = slide1.shapes.add_textbox(Inches(1), Inches(6), Inches(8), Inches(0.5))
    date_frame = date_box.text_frame
    date_frame.text = "2024"
    date_para = date_frame.paragraphs[0]
    date_para.font.size = Pt(20)
    date_para.font.color.rgb = RGBColor(255, 255, 255)
    date_para.alignment = PP_ALIGN.CENTER
    
    # 슬라이드 2: 목차
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide2, "목차", PRIMARY_COLOR)
    
    toc_items = [
        "1. 시스템 개요",
        "2. 사용자 역할별 권한",
        "3. 점검자 사용 가이드",
        "4. 관리자 사용 가이드",
        "5. 시스템 관리자 사용 가이드",
        "6. 주요 기능 요약"
    ]
    
    y_position = 2.0
    for item in toc_items:
        text_box = slide2.shapes.add_textbox(Inches(2), Inches(y_position), Inches(6), Inches(0.5))
        text_frame = text_box.text_frame
        text_frame.text = item
        para = text_frame.paragraphs[0]
        para.font.size = Pt(24)
        para.font.color.rgb = DARK_COLOR
        y_position += 0.6
    
    # 슬라이드 3: 시스템 개요
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide3, "시스템 개요", PRIMARY_COLOR)
    
    overview_text = """
    • HVAC 설비 점검 및 관리를 위한 웹 기반 시스템
    
    • PWA(Progressive Web App) 기술 적용
      - 모바일 앱처럼 설치 가능
      - 오프라인에서도 일부 기능 사용 가능
    
    • 주요 기능
      - QR 코드를 통한 빠른 장비 점검
      - 실시간 점검 내역 관리
      - 장비 데이터 통합 관리
      - 엑셀 데이터 가져오기/내보내기
    """
    
    add_bullet_text(slide3, overview_text, y_position=2.0)
    
    # 슬라이드 4: 사용자 역할
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide4, "사용자 역할별 권한", PRIMARY_COLOR)
    
    # 3개의 역할 박스
    roles = [
        {
            "title": "점검자 (Inspector)",
            "color": SUCCESS_COLOR,
            "permissions": [
                "✓ 장비 점검 수행",
                "✓ QR 스캔 점검",
                "✓ 점검 내역 조회",
                "✓ 자신의 점검 기록 확인"
            ]
        },
        {
            "title": "관리자 (Manager)",
            "color": WARNING_COLOR,
            "permissions": [
                "✓ 점검자 권한 포함",
                "✓ 전체 점검 내역 조회",
                "✓ 대시보드 통계 확인",
                "✓ 장비 목록 조회"
            ]
        },
        {
            "title": "시스템 관리자 (Admin)",
            "color": DANGER_COLOR,
            "permissions": [
                "✓ 관리자 권한 포함",
                "✓ 사용자 관리",
                "✓ 장비 등록/수정/삭제",
                "✓ 현장/건물 관리",
                "✓ 데이터 일괄 작업"
            ]
        }
    ]
    
    x_positions = [0.5, 3.5, 6.5]
    for i, role in enumerate(roles):
        add_role_box(slide4, role, Inches(x_positions[i]), Inches(2.2), Inches(3), Inches(4))
    
    # 슬라이드 5: 점검자 - 메인 화면
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide5, "점검자 - 메인 화면", SUCCESS_COLOR)
    
    inspector_main = """
    🏠 메인 대시보드
    
    • 빠른 점검 시작
      - "장비 점검 시작" 버튼 클릭
      - 현장 → 건물 → 장비 선택
    
    • QR 코드 스캔
      - "QR 스캔" 버튼으로 즉시 점검
      - 장비에 부착된 QR 코드 스캔
    
    • 최근 점검 내역
      - 자신이 수행한 최근 점검 기록
      - 점검 일시, 장비 정보, 상태 확인
    """
    
    add_bullet_text(slide5, inspector_main, y_position=2.0)
    
    # 슬라이드 6: 점검자 - 점검 프로세스
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide6, "점검자 - 점검 수행 프로세스", SUCCESS_COLOR)
    
    # 4단계 프로세스
    steps = [
        {"num": "1", "title": "현장 선택", "desc": "점검할 현장 선택"},
        {"num": "2", "title": "건물 선택", "desc": "해당 현장의 건물 선택"},
        {"num": "3", "title": "장비 선택", "desc": "점검할 장비 선택\n(필터, 검색 가능)"},
        {"num": "4", "title": "점검 입력", "desc": "장비 상태 및\n점검 데이터 입력"}
    ]
    
    for i, step in enumerate(steps):
        x = 1.5 + (i * 2)
        add_process_step(slide6, step, Inches(x), Inches(2.5))
    
    # 화살표 추가 (간단한 텍스트로 표현)
    for i in range(3):
        x = 2.7 + (i * 2)
        arrow_box = slide6.shapes.add_textbox(Inches(x), Inches(3.5), Inches(0.5), Inches(0.5))
        arrow_frame = arrow_box.text_frame
        arrow_frame.text = "→"
        arrow_para = arrow_frame.paragraphs[0]
        arrow_para.font.size = Pt(36)
        arrow_para.font.bold = True
        arrow_para.font.color.rgb = PRIMARY_COLOR
        arrow_para.alignment = PP_ALIGN.CENTER
    
    # 슬라이드 7: 점검자 - 점검 입력 상세
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide7, "점검자 - 점검 데이터 입력", SUCCESS_COLOR)
    
    inspection_details = """
    📝 입력 항목
    
    기본 정보 (자동 입력)
    • 점검자명: 로그인한 사용자 이름 (자동)
    • 점검일시: 현재 날짜/시간 (자동)
    • 장비 정보: 선택한 장비 정보 (자동)
    
    점검 데이터 (직접 입력)
    • 장비 상태: 정상 / 주의 / 경고 / 고장
    • 온도 정보: 실내온도, 설정온도
    • 압력: 고압, 저압
    • 전류: R상, S상, T상 (3상)
    • 특이사항: 텍스트 입력
    • 사진: 장비 사진 첨부 (선택)
    """
    
    add_bullet_text(slide7, inspection_details, y_position=2.0, font_size=18)
    
    # 슬라이드 8: 점검자 - QR 스캔 점검
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide8, "점검자 - QR 코드 스캔 점검", SUCCESS_COLOR)
    
    qr_scan = """
    📱 빠른 점검 방법
    
    1. 메인 화면에서 "QR 스캔" 버튼 클릭
    
    2. 카메라 권한 허용
    
    3. 장비에 부착된 QR 코드 스캔
       → 장비 정보 자동 로드
    
    4. 점검 데이터 입력 화면으로 바로 이동
       → 현장/건물/장비 선택 과정 생략
    
    5. 점검 데이터 입력 및 제출
    
    ⚡ 현장에서 가장 빠르고 편리한 점검 방법!
    """
    
    add_bullet_text(slide8, qr_scan, y_position=2.0)
    
    # 슬라이드 9: 관리자 - 추가 기능
    slide9 = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide9, "관리자 - 추가 기능", WARNING_COLOR)
    
    manager_features = """
    📊 대시보드 접근
    
    • 전체 점검 통계
      - 일간/주간/월간 점검 건수
      - 장비 상태별 통계
      - 점검자별 점검 현황
    
    • 장비 관리
      - 전체 장비 목록 조회
      - 장비별 점검 이력 확인
      - 장비 상태 모니터링
    
    • 점검 내역 조회
      - 전체 점검 내역 확인
      - 기간별, 상태별 필터링
      - 엑셀 다운로드
    """
    
    add_bullet_text(slide9, manager_features, y_position=2.0)
    
    # 슬라이드 10: 시스템 관리자 - 사용자 관리
    slide10 = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide10, "시스템 관리자 - 사용자 관리", DANGER_COLOR)
    
    admin_users = """
    👥 사용자 관리 기능
    
    • 사용자 승인
      - 신규 가입 사용자 승인/거부
      - 역할(점검자/관리자) 지정
    
    • 점검자 추가
      - 새 점검자 계정 생성
      - 이름, 이메일, 비밀번호 설정
      - 역할 지정
    
    • 사용자 수정/삭제
      - 기존 사용자 정보 수정
      - 역할 변경
      - 계정 비활성화/삭제
    
    ⚠️ 주의: 사용자 삭제 시 해당 점검 내역은 유지됨
    """
    
    add_bullet_text(slide10, admin_users, y_position=2.0)
    
    # 슬라이드 11: 시스템 관리자 - 장비 관리
    slide11 = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide11, "시스템 관리자 - 장비 관리", DANGER_COLOR)
    
    admin_equipment = """
    🔧 장비 관리 기능
    
    • 장비 등록 (개별)
      - 장비 ID, 종류, 모델명
      - 현장, 건물, 위치 정보
      - 설치일, 용량 등 상세 정보
    
    • 엑셀 데이터 가져오기
      - 대량의 장비 데이터 일괄 등록
      - 엑셀 템플릿 다운로드
      - 중복 ID 자동 체크
    
    • 장비 수정/삭제
      - 기존 장비 정보 수정
      - 체크박스로 선택 삭제
      - 필터링 및 검색 기능
    
    ⚠️ 주의: 장비 삭제 시 점검 내역은 유지됨
    """
    
    add_bullet_text(slide11, admin_equipment, y_position=2.0)
    
    # 슬라이드 12: 시스템 관리자 - 현장/건물 관리
    slide12 = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide12, "시스템 관리자 - 현장/건물 관리", DANGER_COLOR)
    
    admin_sites = """
    🏢 현장 및 건물 관리
    
    현장 관리
    • 현장 추가: 현장명, 주소, 담당자
    • 현장 수정/삭제
    • 현장별 건물 목록 관리
    
    건물 관리
    • 건물 추가: 건물명, 층수
    • 건물 수정/삭제
    • 현장과 건물 연결
    
    계층 구조
    현장 (Site)
      ↓
    건물 (Building)
      ↓
    장비 (Equipment)
    
    ⚡ 체계적인 장비 위치 관리
    """
    
    add_bullet_text(slide12, admin_sites, y_position=2.0)
    
    # 슬라이드 13: 시스템 관리자 - 점검 내역 관리
    slide13 = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide13, "시스템 관리자 - 점검 내역 관리", DANGER_COLOR)
    
    admin_inspections = """
    📋 점검 내역 관리 기능
    
    • 전체 점검 내역 조회
      - 테이블 형식의 리스트
      - 상세한 점검 정보 표시
    
    • 필터링 및 검색
      - 기간 필터: 오늘, 7일, 30일, 3개월
      - 상태 필터: 정상, 주의, 경고, 고장
      - 점검자 검색
    
    • 선택 삭제
      - 체크박스로 개별 선택
      - 일괄 삭제 기능
      - 삭제 전 확인 메시지
    
    • 삭제된 사용자 표시
      - 노란색 배경으로 강조
      - 경고 아이콘 표시
    """
    
    add_bullet_text(slide13, admin_inspections, y_position=2.0, font_size=18)
    
    # 슬라이드 14: 데이터 관리
    slide14 = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide14, "시스템 관리자 - 데이터 관리", DANGER_COLOR)
    
    data_management = """
    💾 데이터 관리 도구
    
    • 엑셀 데이터 가져오기
      - 장비 데이터 대량 등록
      - 템플릿 다운로드 제공
      - 중복 검사 및 오류 체크
    
    • 엑셀 데이터 내보내기
      - 점검 내역 엑셀 다운로드
      - 장비 목록 엑셀 다운로드
      - 필터 적용된 데이터 내보내기
    
    • 데이터 일괄 삭제
      - 선택한 데이터 일괄 삭제
      - 테스트 데이터 정리
      - 복구 불가능 (신중히 사용)
    
    ⚠️ 중요: 삭제 작업은 되돌릴 수 없습니다!
    """
    
    add_bullet_text(slide14, data_management, y_position=2.0)
    
    # 슬라이드 15: 주요 기능 요약 - QR 코드
    slide15 = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide15, "주요 기능 - QR 코드 시스템", PRIMARY_COLOR)
    
    qr_features = """
    📱 QR 코드 기반 점검 시스템
    
    장점
    • 빠른 점검: 3단계 선택 과정 생략
    • 정확한 장비 식별: 스캔으로 오류 방지
    • 현장 편의성: 모바일에서 즉시 사용
    
    사용 방법
    1. 각 장비에 QR 코드 부착
    2. 점검 시 QR 코드 스캔
    3. 자동으로 장비 정보 로드
    4. 점검 데이터 입력
    
    QR 코드 생성
    • 시스템 관리자가 장비 등록 시 자동 생성
    • 장비 상세 페이지에서 출력 가능
    • 장비 ID를 인코딩하여 생성
    """
    
    add_bullet_text(slide15, qr_features, y_position=2.0)
    
    # 슬라이드 16: 주요 기능 요약 - 대시보드
    slide16 = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide16, "주요 기능 - 대시보드 & 통계", PRIMARY_COLOR)
    
    dashboard_features = """
    📊 실시간 모니터링 대시보드
    
    점검 통계
    • 일간/주간/월간 점검 건수
    • 점검 완료율
    • 미점검 장비 알림
    
    장비 상태 분석
    • 정상/주의/경고/고장 비율
    • 상태별 장비 목록
    • 우선 조치 필요 장비 표시
    
    점검자 활동
    • 점검자별 점검 실적
    • 점검 평균 소요 시간
    • 활동 추이 그래프
    
    최근 점검 내역
    • 실시간 점검 내역 표시
    • 점검 상세 정보 바로 확인
    • 장비별 점검 이력 추적
    """
    
    add_bullet_text(slide16, dashboard_features, y_position=2.0, font_size=18)
    
    # 슬라이드 17: 주요 기능 요약 - 데이터 관리
    slide17 = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide17, "주요 기능 - 데이터 통합 관리", PRIMARY_COLOR)
    
    data_features = """
    💾 효율적인 데이터 관리
    
    엑셀 연동
    • 대량 데이터 일괄 등록
    • 점검 내역 엑셀 내보내기
    • 템플릿 제공으로 편리한 입력
    
    필터 및 검색
    • 다중 조건 필터링
    • 실시간 검색
    • 고급 검색 옵션
    
    데이터 무결성
    • 중복 ID 자동 체크
    • 필수 항목 입력 검증
    • 잘못된 형식 오류 감지
    
    히스토리 관리
    • 모든 점검 내역 보관
    • 장비별 점검 이력 추적
    • 데이터 변경 이력 (추후 추가 가능)
    """
    
    add_bullet_text(slide17, data_features, y_position=2.0, font_size=18)
    
    # 슬라이드 18: 모바일 최적화
    slide18 = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide18, "PWA - 모바일 최적화", PRIMARY_COLOR)
    
    mobile_features = """
    📱 Progressive Web App (PWA)
    
    모바일 앱처럼 사용
    • 홈 화면에 아이콘 추가
    • 전체 화면 모드 실행
    • 앱 스토어 설치 불필요
    
    오프라인 지원
    • 캐시를 통한 빠른 로딩
    • 네트워크 끊김 시에도 일부 기능 사용
    • 재연결 시 자동 동기화
    
    반응형 디자인
    • 모바일, 태블릿, PC 모두 최적화
    • 터치 친화적 UI
    • 큰 버튼과 명확한 레이아웃
    
    카메라 연동
    • QR 코드 스캔
    • 점검 사진 촬영 및 업로드
    • 실시간 이미지 미리보기
    """
    
    add_bullet_text(slide18, mobile_features, y_position=2.0)
    
    # 슬라이드 19: 보안 및 권한
    slide19 = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide19, "보안 및 권한 관리", PRIMARY_COLOR)
    
    security_features = """
    🔒 안전한 데이터 관리
    
    사용자 인증
    • Firebase Authentication 기반
    • 이메일/비밀번호 로그인
    • 비밀번호 재설정 기능
    
    역할 기반 접근 제어
    • 점검자 / 관리자 / 시스템 관리자
    • 역할별 메뉴 및 기능 제한
    • 권한 없는 접근 차단
    
    데이터 보안
    • Firestore Security Rules 적용
    • 사용자별 데이터 접근 권한
    • 읽기/쓰기 권한 세분화
    
    백업 및 복구
    • Firebase 자동 백업
    • 데이터 내보내기 기능
    • 삭제 전 확인 절차
    """
    
    add_bullet_text(slide19, security_features, y_position=2.0)
    
    # 슬라이드 20: 시스템 접속 방법
    slide20 = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide20, "시스템 접속 방법", PRIMARY_COLOR)
    
    access_info = """
    🌐 웹 브라우저 접속
    
    URL: https://noyorc.github.io/hvac-management/
    
    권장 브라우저
    • Google Chrome (권장)
    • Microsoft Edge
    • Safari (iOS/Mac)
    • Samsung Internet (Android)
    
    모바일 설치 방법
    1. 위 URL을 모바일 브라우저에서 접속
    2. 브라우저 메뉴에서 "홈 화면에 추가" 선택
    3. 아이콘이 홈 화면에 생성됨
    4. 앱처럼 실행 가능
    
    초기 계정
    • 신규 가입 후 시스템 관리자 승인 필요
    • 또는 관리자가 직접 계정 생성
    """
    
    add_bullet_text(slide20, access_info, y_position=2.0)
    
    # 슬라이드 21: 사용 시나리오 - 일반 점검
    slide21 = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide21, "사용 시나리오 1: 일반 점검", SUCCESS_COLOR)
    
    scenario1 = """
    📋 정기 점검 시나리오
    
    상황: 점검자 김철수가 본사 A동 장비들을 점검
    
    1. 시스템 로그인
       → 점검자 계정으로 로그인
    
    2. 메인 화면에서 "장비 점검 시작" 클릭
    
    3. 현장 선택: "본사" 선택
    
    4. 건물 선택: "A동" 선택
    
    5. 장비 선택
       → 필터로 "히트펌프" 선택
       → 점검할 장비 클릭
    
    6. 점검 데이터 입력
       → 상태, 온도, 압력 등 입력
       → 특이사항 기록
    
    7. 제출 → 점검 완료!
    """
    
    add_bullet_text(slide21, scenario1, y_position=2.0, font_size=18)
    
    # 슬라이드 22: 사용 시나리오 - QR 점검
    slide22 = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide22, "사용 시나리오 2: QR 스캔 점검", SUCCESS_COLOR)
    
    scenario2 = """
    📱 긴급 점검 시나리오
    
    상황: 고장 신고를 받고 즉시 현장 출동
    
    1. 현장 도착 후 시스템 로그인
       → 모바일로 접속
    
    2. 메인 화면에서 "QR 스캔" 클릭
    
    3. 고장난 장비의 QR 코드 스캔
       → 자동으로 장비 정보 로드
       → 현장/건물 선택 생략!
    
    4. 점검 화면 바로 표시
       → 장비 상태: "고장" 선택
       → 고장 내용 상세 입력
       → 장비 사진 촬영 첨부
    
    5. 제출 → 관리자에게 즉시 알림
    
    ⚡ 총 소요 시간: 2-3분 이내!
    """
    
    add_bullet_text(slide22, scenario2, y_position=2.0, font_size=18)
    
    # 슬라이드 23: 사용 시나리오 - 데이터 분석
    slide23 = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide23, "사용 시나리오 3: 관리자 데이터 분석", WARNING_COLOR)
    
    scenario3 = """
    📊 월간 보고서 작성 시나리오
    
    상황: 관리자가 월간 점검 보고서 작성
    
    1. 대시보드 접속
       → "대시보드" 메뉴 클릭
    
    2. 기간 설정: "최근 30일" 선택
    
    3. 통계 확인
       → 총 점검 건수: 127건
       → 정상: 98건 (77%)
       → 주의/경고: 23건 (18%)
       → 고장: 6건 (5%)
    
    4. 문제 장비 확인
       → "고장" 필터 적용
       → 6건의 고장 장비 상세 확인
       → 조치 이력 검토
    
    5. 데이터 내보내기
       → "엑셀 다운로드" 클릭
       → 보고서에 첨부
    """
    
    add_bullet_text(slide23, scenario3, y_position=2.0, font_size=18)
    
    # 슬라이드 24: 사용 시나리오 - 장비 등록
    slide24 = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide24, "사용 시나리오 4: 신규 장비 등록", DANGER_COLOR)
    
    scenario4 = """
    🔧 신규 현장 장비 등록 시나리오
    
    상황: 신규 지점에 장비 50대 설치 완료
    
    1. 시스템 관리자 로그인
    
    2. 현장/건물 먼저 등록
       → "시스템 관리" → "현장/건물 관리"
       → 현장 추가: "신규 지점"
       → 건물 추가: "본관", "별관"
    
    3. 장비 데이터 준비
       → 엑셀 템플릿 다운로드
       → 장비 정보 50대 입력
       → ID, 종류, 모델, 위치 등
    
    4. 엑셀 데이터 가져오기
       → "장비 관리" → "엑셀 데이터 가져오기"
       → 파일 선택 및 업로드
       → 중복 체크 및 오류 확인
    
    5. 등록 완료
       → 50대 장비 일괄 등록 완료
       → QR 코드 자동 생성
       → 출력하여 장비에 부착
    """
    
    add_bullet_text(slide24, scenario4, y_position=2.0, font_size=16)
    
    # 슬라이드 25: 주요 장점
    slide25 = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide25, "시스템 주요 장점", PRIMARY_COLOR)
    
    # 4개의 장점 박스
    advantages = [
        {
            "icon": "⚡",
            "title": "빠른 점검",
            "desc": "QR 스캔으로\n즉시 점검 시작"
        },
        {
            "icon": "📱",
            "title": "모바일 최적화",
            "desc": "현장에서\n바로 사용"
        },
        {
            "icon": "📊",
            "title": "실시간 통계",
            "desc": "점검 현황\n즉시 파악"
        },
        {
            "icon": "🔒",
            "title": "안전한 관리",
            "desc": "역할별 권한\n데이터 보안"
        }
    ]
    
    for i, adv in enumerate(advantages):
        x = 1.5 + (i % 2) * 4
        y = 2.5 + (i // 2) * 2.5
        add_advantage_box(slide25, adv, Inches(x), Inches(y), Inches(3.5), Inches(2))
    
    # 슬라이드 26: 추후 개선 사항
    slide26 = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide26, "향후 개선 계획", PRIMARY_COLOR)
    
    future_plans = """
    🚀 추가 개발 예정 기능
    
    단기 계획 (1-3개월)
    • 알림 시스템: 미점검 장비 알림, 고장 알림
    • 보고서 자동 생성: PDF 보고서 자동 작성
    • 통계 그래프: 차트 및 시각화 강화
    
    중기 계획 (3-6개월)
    • AI 분석: 장비 고장 예측
    • 일정 관리: 점검 스케줄 자동 생성
    • 다국어 지원: 영어/중국어 등
    
    장기 계획 (6개월 이상)
    • IoT 연동: 센서 데이터 자동 수집
    • 원격 모니터링: 실시간 장비 상태 감시
    • 예지 보전: 고장 전 사전 대응
    """
    
    add_bullet_text(slide26, future_plans, y_position=2.0, font_size=18)
    
    # 슬라이드 27: Q&A / 지원
    slide27 = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide27, "지원 및 문의", PRIMARY_COLOR)
    
    support_info = """
    💬 기술 지원
    
    시스템 문의
    • 사용 방법 문의
    • 오류 신고
    • 기능 개선 제안
    
    교육 및 교육
    • 사용자 교육 자료 제공
    • 동영상 가이드 (준비 중)
    • 1:1 교육 지원
    
    긴급 지원
    • 시스템 장애 시 즉시 대응
    • 데이터 복구 지원
    • 24시간 모니터링
    
    업데이트
    • 정기 업데이트 (월 1회)
    • 보안 패치 즉시 적용
    • 새 기능 추가 시 공지
    """
    
    add_bullet_text(slide27, support_info, y_position=2.0)
    
    # 슬라이드 28: 마무리
    slide28 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 배경색
    background = slide28.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = PRIMARY_COLOR
    
    # 감사 메시지
    thank_you_box = slide28.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
    thank_you_frame = thank_you_box.text_frame
    thank_you_frame.text = "감사합니다"
    thank_you_para = thank_you_frame.paragraphs[0]
    thank_you_para.font.size = Pt(60)
    thank_you_para.font.bold = True
    thank_you_para.font.color.rgb = RGBColor(255, 255, 255)
    thank_you_para.alignment = PP_ALIGN.CENTER
    
    # 부가 메시지
    sub_box = slide28.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1))
    sub_frame = sub_box.text_frame
    sub_frame.text = "HVAC 관리 시스템\n효율적인 설비 관리의 시작"
    for para in sub_frame.paragraphs:
        para.font.size = Pt(24)
        para.font.color.rgb = RGBColor(255, 255, 255)
        para.alignment = PP_ALIGN.CENTER
    
    # 저장
    prs.save('/home/user/webapp/HVAC_Management_System_Guide.pptx')
    print("✅ 프레젠테이션이 생성되었습니다!")
    print("📄 파일명: HVAC_Management_System_Guide.pptx")
    print("📍 위치: /home/user/webapp/")

def add_header(slide, title, color):
    """슬라이드 헤더 추가"""
    header_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    header_frame = header_box.text_frame
    header_frame.text = title
    header_para = header_frame.paragraphs[0]
    header_para.font.size = Pt(36)
    header_para.font.bold = True
    header_para.font.color.rgb = color
    
    # 구분선
    line = slide.shapes.add_shape(
        1,  # 직선
        Inches(0.5), Inches(1.4), Inches(9), Inches(0)
    )
    line.line.color.rgb = color
    line.line.width = Pt(3)

def add_bullet_text(slide, text, y_position=2.0, font_size=20):
    """불릿 포인트 텍스트 추가"""
    text_box = slide.shapes.add_textbox(Inches(1), Inches(y_position), Inches(8), Inches(5))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    text_frame.text = text.strip()
    
    for para in text_frame.paragraphs:
        para.font.size = Pt(font_size)
        para.font.color.rgb = RGBColor(51, 51, 51)
        para.space_after = Pt(8)

def add_role_box(slide, role_data, x, y, width, height):
    """역할 박스 추가"""
    # 박스
    box = slide.shapes.add_shape(
        1,  # 사각형
        x, y, width, height
    )
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(248, 249, 250)
    box.line.color.rgb = role_data["color"]
    box.line.width = Pt(3)
    
    # 제목
    title_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.2), width - Inches(0.4), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = role_data["title"]
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(20)
    title_para.font.bold = True
    title_para.font.color.rgb = role_data["color"]
    title_para.alignment = PP_ALIGN.CENTER
    
    # 권한 목록
    perm_box = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(1), width - Inches(0.6), height - Inches(1.2))
    perm_frame = perm_box.text_frame
    perm_frame.text = "\n".join(role_data["permissions"])
    for para in perm_frame.paragraphs:
        para.font.size = Pt(14)
        para.font.color.rgb = RGBColor(51, 51, 51)
        para.space_after = Pt(6)

def add_process_step(slide, step_data, x, y):
    """프로세스 단계 추가"""
    # 번호 원
    circle = slide.shapes.add_shape(
        9,  # 타원
        x, y, Inches(0.8), Inches(0.8)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(103, 126, 234)
    circle.line.color.rgb = RGBColor(103, 126, 234)
    
    # 번호 텍스트
    num_box = slide.shapes.add_textbox(x, y, Inches(0.8), Inches(0.8))
    num_frame = num_box.text_frame
    num_frame.text = step_data["num"]
    num_para = num_frame.paragraphs[0]
    num_para.font.size = Pt(32)
    num_para.font.bold = True
    num_para.font.color.rgb = RGBColor(255, 255, 255)
    num_para.alignment = PP_ALIGN.CENTER
    num_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # 제목
    title_box = slide.shapes.add_textbox(x - Inches(0.2), y + Inches(1), Inches(1.2), Inches(0.4))
    title_frame = title_box.text_frame
    title_frame.text = step_data["title"]
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(14)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(51, 51, 51)
    title_para.alignment = PP_ALIGN.CENTER
    
    # 설명
    desc_box = slide.shapes.add_textbox(x - Inches(0.3), y + Inches(1.5), Inches(1.4), Inches(0.8))
    desc_frame = desc_box.text_frame
    desc_frame.text = step_data["desc"]
    desc_frame.word_wrap = True
    for para in desc_frame.paragraphs:
        para.font.size = Pt(11)
        para.font.color.rgb = RGBColor(102, 102, 102)
        para.alignment = PP_ALIGN.CENTER

def add_advantage_box(slide, adv_data, x, y, width, height):
    """장점 박스 추가"""
    # 박스
    box = slide.shapes.add_shape(
        1,  # 사각형
        x, y, width, height
    )
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(248, 249, 250)
    box.line.color.rgb = RGBColor(103, 126, 234)
    box.line.width = Pt(2)
    box.shadow.inherit = False
    
    # 아이콘
    icon_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.2), width - Inches(0.4), Inches(0.6))
    icon_frame = icon_box.text_frame
    icon_frame.text = adv_data["icon"]
    icon_para = icon_frame.paragraphs[0]
    icon_para.font.size = Pt(48)
    icon_para.alignment = PP_ALIGN.CENTER
    
    # 제목
    title_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.9), width - Inches(0.4), Inches(0.4))
    title_frame = title_box.text_frame
    title_frame.text = adv_data["title"]
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(18)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(51, 51, 51)
    title_para.alignment = PP_ALIGN.CENTER
    
    # 설명
    desc_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(1.3), width - Inches(0.4), Inches(0.6))
    desc_frame = desc_box.text_frame
    desc_frame.text = adv_data["desc"]
    desc_frame.word_wrap = True
    for para in desc_frame.paragraphs:
        para.font.size = Pt(14)
        para.font.color.rgb = RGBColor(102, 102, 102)
        para.alignment = PP_ALIGN.CENTER

if __name__ == "__main__":
    create_hvac_presentation()
